import streamlit as st
import pdfplumber
import csv
import json
import re
import io
import zipfile
from pathlib import Path
from typing import Optional

st.set_page_config(page_title="문서 → TXT 변환기", page_icon="📄", layout="centered")

st.title("📄 문서 → TXT 변환기")
st.caption("PDF, PPTX, DOCX, XLSX, CSV 파일을 텍스트로 변환합니다")
st.warning("""⚠️ **보안 안내 및 이용 조건**

· 기밀문서·개인정보가 포함된 파일은 업로드하지 마세요.

· 업로드된 파일은 변환 목적으로만 사용되며, 당 서비스는 별도로 저장하지 않습니다.

· 본 서비스는 외부 클라우드 서버(Streamlit)를 경유합니다.

---

*본 서비스 이용 중 발생하는 정보 유출, 데이터 손실 등 모든 문제에 대한 법적 책임은 이용자 본인에게 있으며, 서비스 제공자는 이에 대해 어떠한 책임도 지지 않습니다. 파일 업로드 시 위 내용에 동의한 것으로 간주합니다.*""")

def clean_text(text):
    text = re.sub(r'\n{3,}', '\n\n', text)
    lines = [line.rstrip() for line in text.splitlines()]
    return '\n'.join(lines).strip()

def convert_pdf(file):
    pages = []
    with pdfplumber.open(file) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            tables = page.extract_tables()
            table_text = ""
            for table in tables:
                for row in table:
                    row_clean = [str(cell or "").strip() for cell in row]
                    table_text += " | ".join(row_clean) + "\n"
            combined = text
            if table_text:
                combined += "\n[표]\n" + table_text
            if combined.strip():
                pages.append(f"[페이지 {i}]\n{combined.strip()}")
    return "\n\n".join(pages)

def convert_pptx(file):
    from pptx import Presentation
    prs = Presentation(file)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        parts.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    parts.append(" | ".join(cells))
        if parts:
            slides.append(f"[슬라이드 {i}]\n" + "\n".join(parts))
    return "\n\n".join(slides)

def convert_docx(file):
    from docx import Document
    doc = Document(file)
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            if para.style.name.startswith('Heading'):
                parts.append(f"\n## {text}")
            else:
                parts.append(text)
    for table in doc.tables:
        parts.append("\n[표]")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            parts.append(" | ".join(cells))
    return "\n".join(parts)

def convert_xlsx(file):
    import openpyxl
    wb = openpyxl.load_workbook(file, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"[시트: {sheet_name}]\n" + "\n".join(rows))
    return "\n\n".join(sheets)

def convert_csv(file):
    for encoding in ['utf-8', 'utf-8-sig', 'cp949', 'euc-kr']:
        try:
            content = file.read().decode(encoding)
            file.seek(0)
            reader = csv.reader(content.splitlines())
            rows = []
            for row in reader:
                if any(cell.strip() for cell in row):
                    rows.append(" | ".join(cell.strip() for cell in row))
            return "\n".join(rows)
        except Exception:
            file.seek(0)
            continue
    return ""

def convert_file(uploaded_file):
    ext = Path(uploaded_file.name).suffix.lower()
    if ext == '.pdf':
        return convert_pdf(uploaded_file)
    elif ext in ('.pptx', '.ppt'):
        return convert_pptx(uploaded_file)
    elif ext in ('.docx', '.doc'):
        return convert_docx(uploaded_file)
    elif ext in ('.xlsx', '.xls'):
        return convert_xlsx(uploaded_file)
    elif ext == '.csv':
        return convert_csv(uploaded_file)
    else:
        return None

# 파일 업로드
uploaded_files = st.file_uploader(
    "파일을 여기에 드래그하거나 클릭해서 선택하세요",
    type=['pdf', 'pptx', 'ppt', 'docx', 'doc', 'xlsx', 'xls', 'csv'],
    accept_multiple_files=True
)

if uploaded_files:
    st.divider()
    results = []

    for uploaded_file in uploaded_files:
        with st.spinner(f"{uploaded_file.name} 변환 중..."):
            try:
                text = convert_file(uploaded_file)
                if text:
                    text = clean_text(text)
                    results.append((uploaded_file.name, text))
                    st.success(f"✅ {uploaded_file.name} → 완료 ({len(text):,}자)")
                else:
                    st.warning(f"⚠️ {uploaded_file.name} → 지원하지 않는 형식")
            except Exception as e:
                st.error(f"❌ {uploaded_file.name} → 실패: {str(e)}")

    if results:
        st.divider()

        # 파일 1개면 바로 다운로드
        if len(results) == 1:
            name, text = results[0]
            txt_name = Path(name).stem + ".txt"
            st.download_button(
                label=f"⬇️ {txt_name} 다운로드",
                data=text.encode('utf-8'),
                file_name=txt_name,
                mime='text/plain',
                use_container_width=True
            )

        # 파일 여러 개면 ZIP으로 묶어서 다운로드
        else:
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zf:
                for name, text in results:
                    txt_name = Path(name).stem + ".txt"
                    zf.writestr(txt_name, text.encode('utf-8'))
            zip_buffer.seek(0)

            st.download_button(
                label=f"⬇️ 전체 {len(results)}개 파일 ZIP으로 다운로드",
                data=zip_buffer,
                file_name="converted_txt.zip",
                mime='application/zip',
                use_container_width=True
            )

            st.divider()
            st.caption("개별 파일 다운로드")
            for name, text in results:
                txt_name = Path(name).stem + ".txt"
                st.download_button(
                    label=f"⬇️ {txt_name}",
                    data=text.encode('utf-8'),
                    file_name=txt_name,
                    mime='text/plain',
                    use_container_width=True
                )

st.divider()
st.caption("지원 형식: PDF · PPTX · DOCX · XLSX · CSV  |  무료 · 광고 없음 · 업로드한 파일은 저장되지 않습니다")
